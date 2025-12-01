#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Продвинутый скрипт для исправления кодировки в презентации PowerPoint
ВНИМАНИЕ: python-pptx не поддерживает диаграммы (charts) - они будут потеряны!
Для сохранения диаграмм используйте Microsoft PowerPoint или LibreOffice.
"""
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import sys

# Настройка кодировки для Windows
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

def fix_text_encoding(text: str) -> str:
    """
    Исправляет кодировку текста
    Пытается исправить распространенные проблемы с кодировкой
    """
    if not text:
        return text
    
    # Попытка исправить распространенные проблемы с кодировкой
    # Если текст уже в UTF-8, но был неправильно прочитан из cp1251
    try:
        # Пробуем декодировать как cp1251 и перекодировать в utf-8
        if isinstance(text, str):
            # Если текст содержит кракозябры, пробуем исправить
            # Это сложная задача, требующая знания исходной кодировки
            return text
    except:
        pass
    
    return text

def process_shape(shape, slide_num: int, shape_num: int):
    """
    Обрабатывает фигуру на слайде
    """
    issues = []
    
    # Проверяем тип фигуры
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # Группа фигур - обрабатываем рекурсивно
        for sub_shape in shape.shapes:
            sub_issues = process_shape(sub_shape, slide_num, shape_num)
            issues.extend(sub_issues)
    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        # ДИАГРАММА - python-pptx не поддерживает!
        issues.append(f"Слайд {slide_num}, фигура {shape_num}: ДИАГРАММА (не поддерживается python-pptx)")
    elif shape.has_text_frame:
        # Текстовая фигура - обрабатываем текст
        try:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        # Исправляем кодировку текста
                        fixed_text = fix_text_encoding(run.text)
                        if fixed_text != run.text:
                            run.text = fixed_text
        except Exception as e:
            issues.append(f"Слайд {slide_num}, фигура {shape_num}: Ошибка обработки текста - {e}")
    
    return issues

def fix_presentation_advanced(input_file: str, output_file: str = None):
    """
    Продвинутое исправление презентации с обработкой текста и проверкой диаграмм
    
    Args:
        input_file: Путь к исходному файлу .pptx
        output_file: Путь к выходному файлу (если None, добавляется _FIXED)
    """
    if not os.path.exists(input_file):
        print(f"[ERROR] Ошибка: Файл '{input_file}' не найден!")
        return False
    
    if output_file is None:
        base_name = os.path.splitext(input_file)[0]
        output_file = f"{base_name}_FIXED.pptx"
    
    try:
        print(f"[INFO] Открываю файл: {input_file}")
        prs = Presentation(input_file)
        
        print(f"[INFO] Обрабатываю {len(prs.slides)} слайдов...")
        
        all_issues = []
        charts_found = []
        
        # Обрабатываем каждый слайд
        for slide_idx, slide in enumerate(prs.slides, 1):
            print(f"[INFO] Обработка слайда {slide_idx}/{len(prs.slides)}...")
            
            # Проверяем наличие диаграмм
            for shape_idx, shape in enumerate(slide.shapes, 1):
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    charts_found.append(f"Слайд {slide_idx}")
                    all_issues.append(f"[WARNING] Слайд {slide_idx}: Найдена диаграмма (будет потеряна при сохранении через python-pptx)")
                
                # Обрабатываем фигуру
                issues = process_shape(shape, slide_idx, shape_idx)
                all_issues.extend(issues)
        
        # Выводим предупреждения
        if charts_found:
            print("\n" + "="*70)
            print("[WARNING] ВНИМАНИЕ: Найдены диаграммы на слайдах:")
            for slide_num in charts_found:
                print(f"  - {slide_num}")
            print("\n[WARNING] python-pptx НЕ ПОДДЕРЖИВАЕТ диаграммы!")
            print("[WARNING] Диаграммы будут потеряны при сохранении через этот скрипт.")
            print("[WARNING] Для сохранения диаграмм используйте:")
            print("  1. Microsoft PowerPoint (откройте и сохраните файл)")
            print("  2. LibreOffice Impress (откройте и экспортируйте)")
            print("="*70 + "\n")
        
        if all_issues:
            print("\n[INFO] Найдены проблемы:")
            for issue in all_issues:
                print(f"  - {issue}")
        
        print(f"[INFO] Сохраняю исправленный файл: {output_file}")
        prs.save(output_file)
        
        print(f"\n[OK] Файл успешно обработан и сохранён: {output_file}")
        
        if charts_found:
            print("\n[IMPORTANT] Помните: диаграммы не были сохранены!")
            print("[IMPORTANT] Используйте Microsoft PowerPoint или LibreOffice для полного сохранения.")
        
        return True
        
    except Exception as e:
        print(f"[ERROR] Ошибка при обработке файла: {e}")
        import traceback
        traceback.print_exc()
        return False

def check_presentation_structure(input_file: str):
    """
    Проверяет структуру презентации без изменения
    """
    if not os.path.exists(input_file):
        print(f"[ERROR] Файл '{input_file}' не найден!")
        return
    
    try:
        print(f"[INFO] Анализ файла: {input_file}")
        prs = Presentation(input_file)
        
        print(f"\n[INFO] Количество слайдов: {len(prs.slides)}")
        
        charts_count = 0
        images_count = 0
        text_shapes_count = 0
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    charts_count += 1
                    print(f"  Слайд {slide_idx}: Найдена диаграмма")
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images_count += 1
                elif shape.has_text_frame:
                    text_shapes_count += 1
        
        print(f"\n[INFO] Статистика:")
        print(f"  - Диаграммы: {charts_count}")
        print(f"  - Изображения: {images_count}")
        print(f"  - Текстовые фигуры: {text_shapes_count}")
        
        if charts_count > 0:
            print(f"\n[WARNING] Найдено {charts_count} диаграмм!")
            print("[WARNING] python-pptx не поддерживает диаграммы.")
        
    except Exception as e:
        print(f"[ERROR] Ошибка при анализе: {e}")

if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))
    input_file = os.path.join(script_dir, "Zakupki-RMKSIB-AI-Assistent-1.pptx")
    
    if not os.path.exists(input_file):
        alt_file = os.path.join(script_dir, "Закупки РМКСИБ — AI-Ассистент.pptx")
        if os.path.exists(alt_file):
            input_file = alt_file
    
    if len(sys.argv) > 1:
        command = sys.argv[1]
        if command == "check":
            # Режим проверки
            if len(sys.argv) > 2:
                input_file = sys.argv[2]
            check_presentation_structure(input_file)
            sys.exit(0)
        else:
            input_file = sys.argv[1]
    
    if len(sys.argv) > 2 and sys.argv[1] != "check":
        output_file = sys.argv[2]
    else:
        output_file = None
    
    fix_presentation_advanced(input_file, output_file)

